"""Herestuck Labs build-day deck. Voice follows house-voice.md (voice rules only):
plain noun-phrase headings, no colon-subtitle formula, no em-dashes anywhere, no
'X, not Y' rhetorical pairs, no hedge words, numbers carry their basis inline."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from PIL import Image

D = Path(__file__).resolve().parent.parent / "docs/assets"   # written by render_views.py
OUT = Path(__file__).resolve().parent.parent / "herestuck_boat_build_day.pptx"

INK   = RGBColor(0x14, 0x18, 0x1D)
PAPER = RGBColor(0xF7, 0xF4, 0xEF)
ACC   = RGBColor(0xE4, 0x57, 0x2E)
BLUE  = RGBColor(0x3F, 0x6F, 0xA3)
TEAL  = RGBColor(0x3F, 0xA3, 0x9B)
GOLD  = RGBColor(0xE0, 0xA5, 0x3A)
MUTED = RGBColor(0x6E, 0x76, 0x81)
DIM   = RGBColor(0x9A, 0xA2, 0xAB)

DISP, MONO = "Helvetica Neue", "Menlo"
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)

def slide(dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = INK if dark else PAPER
    return s

def tb(s, x, y, w, h, text, size=18, bold=False, color=INK, font=DISP,
       align=PP_ALIGN.LEFT, space=0, caps=False, line=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln.upper() if caps else ln
        p.alignment = align
        if line: p.line_spacing = line
        for r in p.runs:
            r.font.size, r.font.bold, r.font.name = Pt(size), bold, font
            r.font.color.rgb = color
            if space: r.font._rPr.set("spc", str(int(space * 100)))
    return box

def rule(s, x, y, w, color=ACC, h=0.045):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def pic(s, name, x, y, w=None, h=None):
    p = D / name; iw, ih = Image.open(p).size
    if w and not h: h = w * ih / iw
    if h and not w: w = h * iw / ih
    return s.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))

def head(s, kicker, title, dark=False):
    rule(s, 0.9, 0.72, 0.62)
    tb(s, 0.9, 0.95, 10, 0.3, kicker, 10.5, True, ACC, DISP, caps=True, space=1.8)
    tb(s, 0.9, 1.34, 11.6, 1.0, title, 40, True, PAPER if dark else INK, DISP, line=0.95)

def foot(s, n, dark=False):
    tb(s, 0.9, 6.86, 8, 0.3, "HERESTUCK LABS   ·   CBA BUILD DAY", 8.5, False,
       DIM if dark else MUTED, MONO, caps=True, space=1.2)
    tb(s, 11.3, 6.86, 1.13, 0.3, f"{n:02d}", 8.5, True, DIM if dark else MUTED, MONO,
       align=PP_ALIGN.RIGHT)

def stat(s, x, y, big, cap, color=ACC, size=46, w=3.2, cap_y=None):
    tb(s, x, y, w, 0.85, big, size, True, color, DISP, line=0.9)
    tb(s, x, cap_y if cap_y is not None else y + size / 62.0 + 0.14, w, 0.95,
       cap, 11.5, False, MUTED, DISP, line=1.25)

def bullets(s, x, y, w, items, size=14.5, color=INK, lead=0.24):
    """Advance by the height the subtext actually occupies. A fixed gap works until a
    sub wraps to two lines, and then it silently prints the next title on top of it."""
    yy = y
    for b, sub in items:
        rule(s, x, yy + 0.115, 0.16, ACC, 0.035)
        tb(s, x + 0.34, yy, w, 0.34, b, size, True, color, DISP)
        n = 0
        if sub:
            tb(s, x + 0.34, yy + 0.245, w, 1.2, sub, 11.5, False, MUTED, DISP, line=1.2)
            cpl = max(20, int((w - 0.34) * 12.5))       # chars per line at 11.5 pt
            n = max(1, -(-len(sub) // cpl))
        yy += 0.245 + n * 0.195 + lead
    return yy

def table(s, x, y, w, cols, rows, widths, size=11.5, dark=False):
    n_r, n_c = len(rows) + 1, len(cols)
    shp = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(0.34 * n_r))
    t = shp.table
    t.first_row = False
    for j, cw in enumerate(widths):
        t.columns[j].width = Emu(int(Inches(w) * cw / sum(widths)))
    for j, c in enumerate(cols):
        cell = t.cell(0, j); cell.text = c
        cell.fill.solid(); cell.fill.fore_color.rgb = INK if not dark else RGBColor(0x24,0x2A,0x31)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size, r.font.bold, r.font.name = Pt(9.5), True, MONO
            r.font.color.rgb = PAPER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows, 1):
        for j, v in enumerate(row):
            cell = t.cell(i, j); cell.text = str(v)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if i % 2 else RGBColor(0xEE, 0xEA, 0xE3)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size, r.font.name = Pt(size), MONO
                r.font.color.rgb = INK
                r.font.bold = (j == 0)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return t

# ---------------------------------------------------------------- 01 title
s = slide(dark=True)
rule(s, 0.9, 1.35, 1.05, ACC, 0.06)
tb(s, 0.9, 1.72, 10, 0.35, "Herestuck Labs   ·   CBA Build Day", 12.5, True, ACC, DISP,
   caps=True, space=2.2)
tb(s, 0.86, 2.25, 11.8, 2.2, "Today we finish\nthe bow", 74, True, PAPER, DISP, line=0.92)
tb(s, 0.9, 4.62, 6.6, 1.4,
   "Forty-four chunks make the nose of an eight-foot dinghy. Ten are already printed. "
   "Six machines and a working day gets the other thirty-four.",
   16, False, DIM, DISP, line=1.35)
pic(s, "bow_chunks.png", 7.9, 4.0, w=4.7)
tb(s, 0.9, 6.86, 11.4, 0.3,
   "ASA over gyroid, fiberglassed after   ·   rev-3.4   ·   alexwynn-am.github.io/boat-design",
   9, False, RGBColor(0x55,0x5C,0x64), MONO, caps=True, space=1.2)

# ---------------------------------------------------------------- 02 the boat
s = slide(); head(s, "Context", "The boat the bow belongs to")
pic(s, "hero.png", 6.35, 2.9, w=6.3)
bullets(s, 0.9, 2.7, 5.2, [
    ("8 ft 1 in long, 5 ft 1 in beam", "A shallow-V dinghy with a wide transom, sized around an outboard and a beach landing."),
    ("21.5 kg of ASA, one perimeter over gyroid", "The print is the core. The glass is the structure."),
    ("About 53 lb built", "The plywood version it replaces is 120 to 130 lb, and it oil-canned."),
])
tb(s, 0.9, 5.5, 5.4, 1.0,
   "Every number here comes out of the geometry. The whole model regenerates\n"
   "from one command.", 12, False, MUTED, DISP, line=1.3)
foot(s, 2)

# ---------------------------------------------------------------- 03 four pieces
s = slide(); head(s, "First cut", "Why the boat is four pieces")
pic(s, "four_pieces.png", 5.9, 2.6, w=6.75)
tb(s, 0.9, 2.62, 4.7, 1.2,
   "It has to fit in a car. That one constraint set the whole architecture.",
   17, True, INK, DISP, line=1.25)
bullets(s, 0.9, 3.78, 4.6, [
    ("Centre barge, 42 in wide", "A complete watertight boat on its own. This is the load."),
    ("Two wedge pods", "Sealed buoyancy carrying the outer skin and the beam. They travel separately."),
    ("Bow", "A hollow nose with storage. It nests inside the cockpit with 0.7 in to spare."),
])
foot(s, 3)

# ---------------------------------------------------------------- 04 the bow
s = slide(); head(s, "Second cut", "The bow, in 44 chunks")
pic(s, "bow_chunks.png", 6.35, 2.75, w=6.3)
tb(s, 0.9, 2.62, 5.0, 1.2,
   "The printer is a 256 mm cube. The bow is not.",
   17, True, INK, DISP, line=1.25)
tb(s, 0.9, 3.5, 5.1, 1.9,
   "It gets cut on an axis-aligned grid with 8 mm of margin, and every cut face carries "
   "printed ASA dowels, so a chunk can only go back one way.",
   13.5, False, MUTED, DISP, line=1.4)
stat(s, 0.9, 4.9, "232", "printed dowels\nin the bow alone", ACC, 38, 2.2, cap_y=5.55)
stat(s, 3.3, 4.98, "4.0 mm", "pins in 4.2 mm holes,\n6 mm into each side", INK, 26, 2.7, cap_y=5.55)
foot(s, 4)

# ---------------------------------------------------------------- 05 the day
s = slide(); head(s, "The arithmetic", "Six printers finish the bow")
tb(s, 0.9, 2.5, 5.5, 1.3,
   "The longest chunk left is 2.5 hours. The median is 1.1.",
   19, True, INK, DISP, line=1.2)
tb(s, 0.9, 3.5, 5.4, 1.7,
   "That is why this is a build day. Nobody babysits a forty-hour print, and a machine "
   "that dies at hour two costs one chunk out of thirty-four.",
   13.5, False, MUTED, DISP, line=1.4)
stat(s, 0.9, 5.05, "34", "chunks still to print\n(#011 through #044)", ACC, 40, 2.5, cap_y=5.80)
stat(s, 3.6, 5.05, "40 h", "of machine time,\n3.1 kg of ASA", INK, 40, 2.5, cap_y=5.80)
tb(s, 6.9, 2.52, 5.6, 0.3, "Of the 34 remaining bow chunks, how many land", 11, True, INK, DISP)
table(s, 6.9, 2.9, 5.55,
      ["printers", "8 h day", "10 h day"],
      [["4", "28", "32"], ["6", "34", "34"], ["8", "34", "34"], ["12", "34", "34"]],
      [1.4, 1, 1], size=12)
tb(s, 6.9, 4.75, 5.6, 1.0,
   "Six machines is the whole remaining bow inside eight hours. Four machines leaves "
   "six chunks over, and they are the small ones.",
   11.5, False, MUTED, DISP, line=1.3)
rule(s, 6.9, 5.75, 0.62)
tb(s, 6.9, 5.98, 5.6, 0.6, "Print 000_dowel.stl before anything else.", 14, True, INK, DISP)
foot(s, 5)

# ---------------------------------------------------------------- 06 today
s = slide(); head(s, "The plan", "What we are doing today")
items = [
    ("01", "Load the queue, dowels first",
     "000_dowel.stl, then bow chunks in file order. The numbering runs bottom layer up, so a chunk only ever bonds to lower-numbered neighbours that already exist."),
    ("02", "Dry fit as they come off",
     "Do not wait for a full set. A chunk goes to the bench the moment the plate cools, and #001 through #010 are already there waiting."),
    ("03", "Dowel and bond",
     "Abrade, solvent wipe, epoxy. The ASA to epoxy bond is the least forgiving joint on the boat, and it is the one we make by hand."),
    ("04", "Abrade for glass",
     "Everything that will be glassed gets keyed while there are spare hands to do it."),
]
for i, (n, t, sub) in enumerate(items):
    y = 2.55 + i * 1.08
    tb(s, 0.9, y - 0.04, 0.8, 0.5, n, 26, True, RGBColor(0xD8, 0xD3, 0xCA), MONO)
    tb(s, 1.85, y, 5.6, 0.4, t, 17, True, INK, DISP)
    tb(s, 1.85, y + 0.34, 10.4, 0.6, sub, 12.5, False, MUTED, DISP, line=1.3)
foot(s, 6)

# ---------------------------------------------------------------- 07 orientation
s = slide(dark=True); head(s, "Homework", "Which way up", dark=True)
tb(s, 0.9, 2.62, 6.1, 1.5,
   "Standing each chunk on a flat cut face cut the bow's support by 98 percent.",
   22, True, PAPER, DISP, line=1.18)
tb(s, 0.9, 4.05, 5.9, 2.0,
   "The chunks were being exported in the boat's own frame, which lays most of them "
   "curved side down. A hull has no flat face. A chunk cut off a grid has six of them.\n\n"
   "Support would have cost more ASA than the boat.",
   13.5, False, DIM, DISP, line=1.42)
tb(s, 7.4, 2.52, 5.1, 0.3, "The bow, as cut and as it now prints", 11, True, PAPER, DISP)
table(s, 7.4, 2.9, 5.1,
      ["", "as cut", "oriented"],
      [["support", "9,966 cm²", "223 cm²"],
       ["envelope", "72.6 L", "0.25 L"],
       ["print time", "about 200 h", "54 h"],
       ["needs support", "most of 44", "6 of 44"]],
      [1.6, 1.2, 1.1], dark=True)
tb(s, 7.4, 5.4, 5.1, 0.8, "Roughly four times faster, for a rotation.",
   14, True, ACC, DISP)
foot(s, 7, dark=True)

# ---------------------------------------------------------------- 08 joints
s = slide(); head(s, "Assembly", "How the pieces find each other")
bullets(s, 0.9, 2.6, 5.3, [
    ("232 printed ASA dowels in the bow", "4.0 mm pins in 4.2 mm holes, 6 mm into each side of every cut face."),
    ("Vertical dovetail keys", "The four big pieces drop together from above and lock against pulling apart."),
    ("16 bolts, all above the waterline", "They clamp. The keys resist lift. Nothing structural sits below the water."),
])
tb(s, 6.6, 2.6, 5.9, 2.7,
   "Joint clearance is 2.0 mm per face, which sounds loose until you count the glass "
   "going into it. Both mating walls are exterior surfaces, they both take 6 oz cloth, "
   "and that ate the original 0.6 mm on its own.\n\n"
   "Tongues take cloth. Sockets take neat epoxy, because cloth only bridges a 16 mm "
   "flared slot and traps air behind itself.",
   13.5, False, MUTED, DISP, line=1.45)
rule(s, 6.6, 5.6, 0.62)
tb(s, 6.6, 5.83, 5.9, 0.6, "Loose is nearly free. The bolts do the clamping.",
   14, True, INK, DISP)
foot(s, 8)

# ---------------------------------------------------------------- 09 glass
s = slide(); head(s, "After the plastic", "The print is not the boat yet")
tb(s, 0.9, 2.6, 6.2, 1.3,
   "A 0.5 mm wall over 8 percent gyroid weeps through its own layer lines. "
   "The glass is what makes it a boat.",
   19, True, INK, DISP, line=1.22)
bullets(s, 0.9, 4.05, 5.4, [
    ("6 oz cloth everywhere", "One fabric for the whole boat, chosen for wet-out and for drape over printed layer lines."),
    ("Doubled on the bottom, second layer at 45 degrees", "Beaching and slamming loads land there."),
])
tb(s, 7.0, 2.62, 5.5, 3.2,
   "The printed wall carries about the same weight per unit area as the plywood it "
   "replaces, so the thickness is close to free. The mass lives in perimeters and glass.\n\n"
   "Infill is half the boat's mass, perimeters are 30 percent, glass is 20. One percent "
   "of infill is 2.7 lb.\n\n"
   "Fifty pounds is reachable. Forty is not, without a smaller boat.",
   13.5, False, MUTED, DISP, line=1.45)
foot(s, 9)

# ---------------------------------------------------------------- 10 failure
s = slide(); head(s, "Expect this", "What will go wrong")
warn = [
    ("ASA lifts off the plate",
     "Door shut, chamber warm before layer one, 8 mm brim on every chunk. This is the failure you will actually see today."),
    ("Six bow chunks still want support",
     "#043 and #044 at 52 cm², #027 and #028 at 24, #013 and #016 at 10. Use normal supports at threshold 30, with 'on build plate only' and 'support critical regions only' both OFF."),
    ("#034 and #035 stand 13.7:1 tall",
     "They sit on 17 cm² of plate, and that stance is what makes them support-free. Watch the first layer, then leave them alone."),
    ("286 of the bow's 464 dowel bores print horizontally",
     "They bridge their own roof and can droop into a 4.2 mm hole. Dry fit before glue. If they bind, reprint the dowels thinner."),
]
for i, (t, sub) in enumerate(warn):
    y = 2.48 + i * 1.03
    rule(s, 0.9, y + 0.13, 0.22, ACC, 0.05)
    tb(s, 1.35, y, 6.2, 0.4, t, 15.5, True, INK, DISP)
    tb(s, 1.35, y + 0.32, 10.9, 0.7, sub, 12, False, MUTED, DISP, line=1.28)
foot(s, 10)

# ---------------------------------------------------------------- 11 close
s = slide(dark=True)
rule(s, 0.9, 2.35, 1.05, ACC, 0.06)
tb(s, 0.86, 2.72, 11.6, 1.8, "Take a chunk.\nBring it back bonded.", 52, True,
   PAPER, DISP, line=0.98)
tb(s, 0.9, 4.72, 6.8, 1.3,
   "Forty-four chunks, 232 dowels, one bow. manifest.csv says which chunk mates to which, "
   "and PROFILE.txt in the bow folder says how to print all of it.",
   15, False, DIM, DISP, line=1.4)
rule(s, 0.9, 5.95, 0.62)
tb(s, 0.9, 6.16, 9.0, 0.45, "alexwynn-am.github.io/boat-design", 22, True, ACC, DISP)
tb(s, 0.9, 6.86, 10, 0.3,
   "000_DOWEL.STL FIRST   ·   THEN #011 UPWARD   ·   BRIM EVERYTHING",
   9, False, RGBColor(0x55,0x5C,0x64), MONO, caps=True, space=1.2)

prs.save(str(OUT))
print("wrote", OUT, f"{OUT.stat().st_size/1024:.0f} KB")
